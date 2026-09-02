"""Turn source files into chunks with human-readable locators.

Formats:
  .md          split by heading; locator is the heading path ("Intro > Setup")
  .pdf         one chunk per page; locator "page N"
  .pptx        one chunk per slide (including speaker notes); locator "slide N"
  .docx        split by Heading 1 / Heading 2; locator is the heading path
  .xlsx        one chunk per worksheet, formulas kept; locator "sheet Name"
  .html        visible text, script/style dropped; locator is the <title>
  .txt         packed plain text; locator "part N"

Target ~800 tokens per chunk with ~100 tokens of overlap, never splitting
mid-sentence. Tokens are estimated as chars/4 - close enough for packing.

Oversized sections/pages are packed into "(part N)" sub-chunks. A single
"sentence" longer than the chunk target (a wall of text with no punctuation)
is force-split at word boundaries - the one place the no-mid-sentence rule
physically cannot hold.
"""

from __future__ import annotations

import hashlib
import re
from dataclasses import dataclass, field
from pathlib import Path
from typing import Callable

from .errors import ParseError

TARGET_TOKENS = 800
OVERLAP_TOKENS = 100

SUPPORTED_SUFFIXES = {".md", ".pdf", ".pptx", ".docx", ".txt", ".xlsx", ".html", ".htm"}


@dataclass
class Chunk:
    locator: str
    text: str

    @property
    def content_hash(self) -> str:
        return hashlib.sha256(self.text.encode("utf-8")).hexdigest()[:16]


def estimate_tokens(text: str) -> int:
    return max(1, len(text) // 4)


_SENTENCE_RE = re.compile(r"(?<=[.!?])\s+|\n{2,}")


def split_sentences(text: str) -> list[str]:
    """Split into sentence-ish units on sentence enders and blank lines."""
    parts = [p.strip() for p in _SENTENCE_RE.split(text)]
    return [p for p in parts if p]


def _force_split(
    sentence: str, target: int, count_tokens: Callable[[str], int] = estimate_tokens
) -> list[str]:
    """Split an oversized single sentence at word boundaries, measuring the
    accumulated text rather than summing per-word counts (which underestimate,
    since a word costs fewer tokens on its own than in context)."""
    words = sentence.split()
    out: list[str] = []
    cur: list[str] = []
    for w in words:
        if cur and count_tokens(" ".join([*cur, w])) > target:
            out.append(" ".join(cur))
            cur = []
        cur.append(w)
    if cur:
        out.append(" ".join(cur))
    # A single token longer than the target (a giant base64 blob) cannot be
    # split at a word boundary; cut it by characters so it still fits.
    final: list[str] = []
    for piece in out:
        while count_tokens(piece) > target and len(piece) > 4:
            keep = max(4, int(len(piece) * target / max(1, count_tokens(piece))))
            final.append(piece[:keep])
            piece = piece[keep:]
        if piece:
            final.append(piece)
    return final


def pack_sentences(
    text: str,
    target: int = TARGET_TOKENS,
    overlap: int = OVERLAP_TOKENS,
    count_tokens: Callable[[str], int] = estimate_tokens,
) -> list[str]:
    """Pack sentences into ~target-token chunks with ~overlap-token carryover.

    `count_tokens` defaults to the chars/4 estimate, which keeps this module
    model-agnostic and fast to test. The indexer passes the embedding model's
    real tokenizer instead, because the estimate is off by ~1.45x on real
    documents and anything past the model's window is silently discarded.
    """
    sentences: list[str] = []
    for s in split_sentences(text):
        if count_tokens(s) > target:
            sentences.extend(_force_split(s, target, count_tokens))
        else:
            sentences.append(s)
    if not sentences:
        return []

    chunks: list[str] = []
    cur: list[str] = []
    cur_tokens = 0
    for s in sentences:
        st = count_tokens(s)
        if cur and cur_tokens + st > target:
            chunks.append("\n".join(cur))
            # Carry trailing sentences into the next chunk, but never so many
            # that the carry-over plus the incoming sentence blows the target
            # (that is how chunks used to end up multiples of their size).
            carried: list[str] = []
            carried_tokens = 0
            room = max(0, min(overlap, target - st))
            for prev in reversed(cur):
                pt = count_tokens(prev)
                if carried_tokens + pt > room:
                    break
                carried.insert(0, prev)
                carried_tokens += pt
            # Overlap must never be the whole previous chunk, or packing stalls.
            if carried_tokens >= cur_tokens:
                carried, carried_tokens = [], 0
            cur = list(carried)
            cur_tokens = carried_tokens
        cur.append(s)
        cur_tokens += st
    if cur:
        chunks.append("\n".join(cur))
    return chunks


@dataclass
class Packer:
    """How to size chunks. Defaults keep this module model-agnostic; the
    indexer supplies the embedding model's real tokenizer and window."""

    target: int = TARGET_TOKENS
    overlap: int = OVERLAP_TOKENS
    count_tokens: Callable[[str], int] = estimate_tokens

    def pack(self, text: str) -> list[str]:
        return pack_sentences(text, self.target, self.overlap, self.count_tokens)


DEFAULT_PACKER = Packer()


def _packed(locator: str, text: str, packer: Packer = DEFAULT_PACKER) -> list[Chunk]:
    """Pack a section's text; suffix '(part N)' only when it splits."""
    pieces = packer.pack(text)
    if not pieces:
        return []
    if len(pieces) == 1:
        return [Chunk(locator=locator, text=pieces[0])]
    return [
        Chunk(locator=f"{locator} (part {i + 1})", text=p)
        for i, p in enumerate(pieces)
    ]


# ---------------------------------------------------------------- markdown

_HEADING_RE = re.compile(r"^(#{1,6})\s+(.*?)\s*#*\s*$")
_FENCE_RE = re.compile(r"^(```|~~~)")


def chunk_markdown(text: str, *, title: str = "", packer: Packer = DEFAULT_PACKER) -> list[Chunk]:
    """Split by headings, locator = heading path. Fenced code is not scanned
    for headings. Content before the first heading gets the file title (or
    '(preamble)') as its locator."""
    sections: list[tuple[str, list[str]]] = []  # (heading_path, lines)
    path_stack: list[tuple[int, str]] = []      # (level, heading text)
    cur_lines: list[str] = []
    in_fence = False
    fence_marker = ""

    def flush() -> None:
        nonlocal cur_lines
        body = "\n".join(cur_lines).strip()
        if body:
            hp = " > ".join(h for _, h in path_stack) or (title or "(preamble)")
            sections.append((hp, cur_lines))
        cur_lines = []

    for line in text.splitlines():
        fm = _FENCE_RE.match(line.strip())
        if fm:
            if not in_fence:
                in_fence, fence_marker = True, fm.group(1)
            elif line.strip().startswith(fence_marker):
                in_fence = False
            cur_lines.append(line)
            continue
        hm = None if in_fence else _HEADING_RE.match(line)
        if hm:
            flush()
            level = len(hm.group(1))
            heading = hm.group(2).strip()
            while path_stack and path_stack[-1][0] >= level:
                path_stack.pop()
            path_stack.append((level, heading))
            cur_lines.append(line)
        else:
            cur_lines.append(line)
    flush()

    chunks: list[Chunk] = []
    for heading_path, lines in sections:
        chunks.extend(_packed(heading_path, "\n".join(lines).strip(), packer))
    return chunks


# ---------------------------------------------------------------- pdf

def chunk_pdf(path: Path, packer: Packer = DEFAULT_PACKER) -> list[Chunk]:
    import pymupdf

    chunks: list[Chunk] = []
    with pymupdf.open(path) as doc:
        if doc.is_encrypted and not doc.authenticate(""):
            raise ParseError(str(path), "PDF is password-protected")
        for page in doc:
            text = page.get_text().strip()
            if not text:
                continue
            chunks.extend(_packed(f"page {page.number + 1}", text, packer))
    return chunks


# ---------------------------------------------------------------- pptx

def _pptx_shape_text(shape, out: list[str]) -> None:
    """Collect text from a shape, descending into groups and tables. A lecture
    slide that is only a table would otherwise contribute nothing."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for child in shape.shapes:
            _pptx_shape_text(child, out)
        return
    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            cells = [c.text.strip() for c in row.cells]
            if any(cells):
                out.append(" | ".join(cells))
        return
    if shape.has_text_frame:
        t = shape.text_frame.text.strip()
        if t:
            out.append(t)


def chunk_pptx(path: Path, packer: Packer = DEFAULT_PACKER) -> list[Chunk]:
    from pptx import Presentation

    prs = Presentation(str(path))
    chunks: list[Chunk] = []
    for i, slide in enumerate(prs.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            _pptx_shape_text(shape, parts)
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                parts.append(f"[Speaker notes] {notes}")
        text = "\n".join(parts).strip()
        if not text:
            continue
        chunks.extend(_packed(f"slide {i}", text, packer))
    return chunks


# ---------------------------------------------------------------- docx

_W_NS = "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}"


def _chunk_docx_huge(path: Path, packer: Packer = DEFAULT_PACKER) -> list[Chunk]:
    """Fallback for documents lxml refuses under its default limits.

    Scanned readings embed image data as enormous attribute values (one real
    file has an 11 MB attribute inside a 60 MB document.xml), which trips
    lxml's 10 MB attribute cap. Re-parse the body directly with huge_tree
    enabled and read styles from w:pStyle, since python-docx does not expose
    parser options.
    """
    import zipfile

    from lxml import etree

    parser = etree.XMLParser(huge_tree=True, recover=True, resolve_entities=False)
    with zipfile.ZipFile(path) as z:
        with z.open("word/document.xml") as f:
            tree = etree.parse(f, parser)

    sections: list[tuple[str, list[str]]] = []
    h1: str | None = None
    h2: str | None = None
    cur: list[str] = []

    def flush() -> None:
        nonlocal cur
        body = "\n".join(cur).strip()
        if body:
            hp = " > ".join(x for x in (h1, h2) if x) or "(body)"
            sections.append((hp, cur))
        cur = []

    for para in tree.iter(f"{_W_NS}p"):
        text = "".join(t.text or "" for t in para.iter(f"{_W_NS}t")).strip()
        style_el = para.find(f"{_W_NS}pPr/{_W_NS}pStyle")
        style = (style_el.get(f"{_W_NS}val") or "") if style_el is not None else ""
        normalized = style.replace(" ", "").lower()
        if normalized == "heading1":
            flush()
            h1, h2 = text or h1, None
        elif normalized == "heading2":
            flush()
            h2 = text or h2
        if text:
            cur.append(text)
    flush()

    chunks: list[Chunk] = []
    for heading_path, lines in sections:
        chunks.extend(_packed(heading_path, "\n".join(lines).strip(), packer))
    return chunks


def _iter_docx_blocks(parent):
    """Yield a container's Paragraphs and Tables in document order.

    document.paragraphs and document.tables are both shallow: neither
    descends into a table cell. Newsletter-style .docx files wrap their whole
    body in a layout table containing a nested content table, so both lists
    come back empty and the file yields nothing at all.
    """
    from docx.document import Document as _Document
    from docx.oxml.table import CT_Tbl
    from docx.oxml.text.paragraph import CT_P
    from docx.table import Table, _Cell
    from docx.text.paragraph import Paragraph

    if isinstance(parent, _Document):
        parent_elm = parent.element.body
    elif isinstance(parent, _Cell):
        parent_elm = parent._tc
    else:
        raise ValueError(f"Cannot iterate blocks of {type(parent).__name__}")

    for child in parent_elm.iterchildren():
        if isinstance(child, CT_P):
            yield Paragraph(child, parent)
        elif isinstance(child, CT_Tbl):
            yield Table(child, parent)


def _docx_table_lines(table, out: list[str], depth: int = 0) -> None:
    """Flatten a table to 'cell | cell' lines, recursing into nested tables."""
    from docx.table import Table

    if depth > 6:  # pathological nesting guard
        return
    for row in table.rows:
        cells: list[str] = []
        for cell in row.cells:
            texts: list[str] = []
            for block in _iter_docx_blocks(cell):
                if isinstance(block, Table):
                    _docx_table_lines(block, out, depth + 1)
                else:
                    t = block.text.strip()
                    if t:
                        texts.append(t)
            cells.append(" ".join(texts))
        if any(cells):
            out.append(" | ".join(cells))


def chunk_docx(path: Path, packer: Packer = DEFAULT_PACKER) -> list[Chunk]:
    import docx
    from docx.table import Table
    from lxml.etree import XMLSyntaxError

    try:
        document = docx.Document(str(path))
    except XMLSyntaxError:
        return _chunk_docx_huge(path, packer)
    sections: list[tuple[str, list[str]]] = []
    h1: str | None = None
    h2: str | None = None
    cur: list[str] = []

    def flush() -> None:
        nonlocal cur
        body = "\n".join(cur).strip()
        if body:
            hp = " > ".join(x for x in (h1, h2) if x) or "(body)"
            sections.append((hp, cur))
        cur = []

    # One ordered walk over paragraphs AND tables, so table content keeps the
    # heading it appeared under instead of being dumped in a trailing section.
    for block in _iter_docx_blocks(document):
        if isinstance(block, Table):
            _docx_table_lines(block, cur)
            continue
        style = (block.style.name or "") if block.style else ""
        text = block.text.strip()
        if style.startswith("Heading 1"):
            flush()
            h1, h2 = text or h1, None
            if text:
                cur.append(text)
        elif style.startswith("Heading 2"):
            flush()
            h2 = text or h2
            if text:
                cur.append(text)
        elif text:
            cur.append(text)
    flush()

    chunks: list[Chunk] = []
    for heading_path, lines in sections:
        chunks.extend(_packed(heading_path, "\n".join(lines).strip(), packer))
    return chunks


# ---------------------------------------------------------------- xlsx

def _cell_repr(value) -> str:
    """Render a cell. Formulas are kept as written: in a spreadsheet course
    the formula IS the content, and '=PMT(B4/12,B5*12,-B3)' is what a
    question about loan payments should match."""
    if value is None:
        return ""
    if isinstance(value, str):
        return value.strip()
    if isinstance(value, float) and value.is_integer():
        return str(int(value))
    return str(value)


def chunk_xlsx(path: Path, packer: Packer = DEFAULT_PACKER) -> list[Chunk]:
    """One chunk per worksheet (split into parts when large); locator is the
    sheet name. Rows render as 'a | b | c', trailing empties trimmed."""
    import openpyxl

    wb = openpyxl.load_workbook(path, data_only=False, read_only=True)
    try:
        chunks: list[Chunk] = []
        for ws in wb.worksheets:
            lines: list[str] = []
            for row in ws.iter_rows(values_only=True):
                # Empty cells are dropped rather than rendered as separators:
                # spreadsheets are full of layout gaps, and ' | | | ' is pure
                # noise in an embedding.
                cells = [c for c in (_cell_repr(v) for v in row) if c]
                if cells:
                    lines.append(" | ".join(cells))
            text = "\n".join(lines).strip()
            if not text:
                continue
            chunks.extend(_packed(f"sheet {ws.title}", text, packer))
        return chunks
    finally:
        wb.close()


# ---------------------------------------------------------------- html

def chunk_html(path: Path, packer: Packer = DEFAULT_PACKER) -> list[Chunk]:
    """Visible text only. Script/style content is dropped rather than indexed
    as if it were prose."""
    from lxml import html as lxml_html

    raw = _read_text(path)
    if not raw.strip():
        return []
    doc = lxml_html.fromstring(raw)
    for bad in doc.xpath("//script | //style | //noscript"):
        bad.getparent().remove(bad)
    title = ""
    found = doc.xpath("//title/text()")
    if found:
        title = str(found[0]).strip()
    text = "\n".join(
        line.strip() for line in doc.text_content().splitlines() if line.strip()
    )
    if not text:
        return []
    return _packed(title or path.stem, text, packer)


# ---------------------------------------------------------------- txt

def chunk_txt(text: str, packer: Packer = DEFAULT_PACKER) -> list[Chunk]:
    pieces = packer.pack(text)
    if not pieces:
        return []
    if len(pieces) == 1:
        return [Chunk(locator="full text", text=pieces[0])]
    return [Chunk(locator=f"part {i + 1}", text=p) for i, p in enumerate(pieces)]


# ---------------------------------------------------------------- dispatch

def chunk_file(
    path: Path,
    *,
    count_tokens: Callable[[str], int] | None = None,
    target_tokens: int | None = None,
) -> list[Chunk]:
    """Parse one file into chunks. Raises ParseError with the reason on failure.

    Pass the embedding model's tokenizer and window as count_tokens/
    target_tokens so chunks fit what the model actually reads; without them
    the chars/4 estimate and the ~800-token default are used.
    """
    packer = DEFAULT_PACKER
    if count_tokens is not None or target_tokens is not None:
        target = target_tokens or TARGET_TOKENS
        packer = Packer(
            target=target,
            # Keep the spec's overlap:target ratio when the target shrinks to
            # fit the model, rather than overlapping a quarter of each chunk.
            overlap=max(1, round(OVERLAP_TOKENS * target / TARGET_TOKENS)),
            count_tokens=count_tokens or estimate_tokens,
        )
    suffix = path.suffix.lower()
    try:
        if suffix == ".md":
            return chunk_markdown(_read_text(path), title=path.stem, packer=packer)
        if suffix == ".txt":
            return chunk_txt(_read_text(path), packer)
        if suffix == ".pdf":
            return chunk_pdf(path, packer)
        if suffix == ".pptx":
            return chunk_pptx(path, packer)
        if suffix == ".docx":
            return chunk_docx(path, packer)
        if suffix == ".xlsx":
            return chunk_xlsx(path, packer)
        if suffix in (".html", ".htm"):
            return chunk_html(path, packer)
    except ParseError:
        raise
    except Exception as e:
        raise ParseError(str(path), f"{type(e).__name__}: {e}") from e
    raise ParseError(str(path), f"Unsupported file type '{suffix}'")


# Byte-order marks, longest first so UTF-32 is not read as UTF-16.
_BOMS: tuple[tuple[bytes, str], ...] = (
    (b"\xff\xfe\x00\x00", "utf-32-le"),
    (b"\x00\x00\xfe\xff", "utf-32-be"),
    (b"\xef\xbb\xbf", "utf-8-sig"),
    (b"\xff\xfe", "utf-16-le"),
    (b"\xfe\xff", "utf-16-be"),
)


def _read_text(path: Path) -> str:
    """Decode a text file, honoring any BOM.

    Two traps this avoids: a UTF-8 BOM left on the front of a markdown file
    hides its first '#' from the heading regex, collapsing the whole document
    into one locator; and UTF-16 decodes 'successfully' as latin-1, silently
    indexing every character interleaved with NULs.
    """
    raw = path.read_bytes()
    for bom, encoding in _BOMS:
        if raw.startswith(bom):
            return raw.decode(encoding).lstrip("﻿")
    try:
        return raw.decode("utf-8")
    except UnicodeDecodeError:
        pass
    # No BOM but NUL-heavy: UTF-16 written without one. Guessing latin-1 here
    # would produce mojibake that indexes cleanly and retrieves as garbage.
    head = raw[:4096]
    if head.count(0) > len(head) // 4:
        for encoding in ("utf-16-le", "utf-16-be"):
            try:
                return raw.decode(encoding)
            except UnicodeDecodeError:
                continue
        raise ParseError(str(path), "looks like UTF-16 but could not be decoded")
    return raw.decode("latin-1")

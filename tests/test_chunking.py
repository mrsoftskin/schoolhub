"""Chunking per format: heading paths, page/slide locators, overlap,
sentence-boundary preservation."""

from __future__ import annotations

import re

from brain.chunking import (
    OVERLAP_TOKENS,
    TARGET_TOKENS,
    chunk_docx,
    chunk_file,
    chunk_markdown,
    chunk_pdf,
    chunk_pptx,
    chunk_txt,
    estimate_tokens,
    pack_sentences,
)
from brain.errors import ParseError


def sentence(i: int) -> str:
    return f"This is test sentence number {i} with some padding words to give it heft."


# ------------------------------------------------------------- packing

def test_pack_respects_target_and_never_splits_mid_sentence():
    text = " ".join(sentence(i) for i in range(200))
    chunks = pack_sentences(text)
    assert len(chunks) > 1
    for c in chunks:
        assert estimate_tokens(c) <= TARGET_TOKENS + 40  # one-sentence tolerance
        # Every chunk ends exactly at a sentence boundary.
        assert c.rstrip().endswith("heft.")


def test_pack_overlap_carries_trailing_sentences():
    text = " ".join(sentence(i) for i in range(200))
    chunks = pack_sentences(text)
    for a, b in zip(chunks, chunks[1:]):
        a_sents = a.split("\n")
        b_sents = b.split("\n")
        # The next chunk starts with the tail of the previous one.
        carried = 0
        while carried < len(b_sents) and b_sents[carried] in a_sents:
            carried += 1
        assert carried >= 1, "no overlap between consecutive chunks"
        carried_tokens = sum(estimate_tokens(s) for s in b_sents[:carried])
        assert carried_tokens >= OVERLAP_TOKENS * 0.5


def test_pack_handles_wall_of_text_without_punctuation():
    text = "word " * 5000
    chunks = pack_sentences(text)
    assert len(chunks) > 1
    for c in chunks:
        assert estimate_tokens(c) <= TARGET_TOKENS + 40


# ------------------------------------------------------------- markdown

def test_markdown_heading_paths():
    md = "\n".join([
        "intro before any heading.",
        "# Alpha",
        "Alpha body text.",
        "## Beta",
        "Beta body text.",
        "## Gamma",
        "Gamma body text.",
        "# Delta",
        "Delta body text.",
    ])
    chunks = chunk_markdown(md, title="testfile")
    locs = [c.locator for c in chunks]
    assert locs == ["testfile", "Alpha", "Alpha > Beta", "Alpha > Gamma", "Delta"]
    assert "Beta body text." in chunks[2].text


def test_markdown_ignores_headings_inside_code_fences():
    md = "# Real\nbody\n```\n# not a heading\n```\nmore body\n"
    chunks = chunk_markdown(md)
    assert [c.locator for c in chunks] == ["Real"]
    assert "# not a heading" in chunks[0].text


def test_markdown_oversized_section_gets_parts():
    md = "# Big\n" + " ".join(sentence(i) for i in range(300))
    chunks = chunk_markdown(md)
    assert len(chunks) > 1
    assert all(c.locator.startswith("Big (part ") for c in chunks)


# ------------------------------------------------------------- pdf

def test_pdf_page_locators(tmp_path):
    import pymupdf

    doc = pymupdf.open()
    for i in range(3):
        page = doc.new_page()
        page.insert_text((72, 72), f"Content of page number {i + 1}. It has a sentence.")
    path = tmp_path / "t.pdf"
    doc.save(str(path))
    doc.close()

    chunks = chunk_pdf(path)
    assert [c.locator for c in chunks] == ["page 1", "page 2", "page 3"]
    assert "page number 2" in chunks[1].text


# ------------------------------------------------------------- pptx

def test_pptx_slide_locators_and_notes(tmp_path):
    from pptx import Presentation

    prs = Presentation()
    layout = prs.slide_layouts[1]
    for i in range(2):
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = f"Slide title {i + 1}"
        slide.placeholders[1].text = f"Bullet content {i + 1}"
        slide.notes_slide.notes_text_frame.text = f"Note for slide {i + 1}"
    path = tmp_path / "t.pptx"
    prs.save(str(path))

    chunks = chunk_pptx(path)
    assert [c.locator for c in chunks] == ["slide 1", "slide 2"]
    assert "Slide title 2" in chunks[1].text
    assert "Note for slide 2" in chunks[1].text
    assert "[Speaker notes]" in chunks[1].text


# ------------------------------------------------------------- docx

def test_docx_heading_paths(tmp_path):
    import docx

    d = docx.Document()
    d.add_paragraph("Preamble text before headings.")
    d.add_heading("Chapter One", level=1)
    d.add_paragraph("Body of chapter one.")
    d.add_heading("Section A", level=2)
    d.add_paragraph("Body of section A.")
    d.add_heading("Chapter Two", level=1)
    d.add_paragraph("Body of chapter two.")
    path = tmp_path / "t.docx"
    d.save(str(path))

    chunks = chunk_docx(path)
    locs = [c.locator for c in chunks]
    assert locs == ["(body)", "Chapter One", "Chapter One > Section A", "Chapter Two"]
    assert "Body of section A." in chunks[2].text


# ------------------------------------------------------------- txt / dispatch

def test_docx_nested_tables_are_not_lost(tmp_path):
    """Newsletter-style .docx files wrap the body in a layout table holding a
    nested content table. document.paragraphs and document.tables are both
    shallow, so such a file used to index as zero chunks."""
    import docx

    d = docx.Document()
    outer = d.add_table(rows=1, cols=1)
    cell = outer.rows[0].cells[0]
    inner = cell.add_table(rows=2, cols=2)
    inner.rows[0].cells[0].text = "Houston Office Under Contract"
    inner.rows[0].cells[1].text = "Reappraised this quarter"
    inner.rows[1].cells[0].text = "One City Centre"
    inner.rows[1].cells[1].text = "Servicer commentary updated"
    path = tmp_path / "newsletter.docx"
    d.save(str(path))

    # The shallow accessors really do come back empty for this shape.
    reopened = docx.Document(str(path))
    assert reopened.paragraphs == [] or all(not p.text.strip() for p in reopened.paragraphs)
    assert reopened.tables[0].rows[0].cells[0].text.strip() == ""

    chunks = chunk_file(path)
    assert chunks, "nested table content must not be lost"
    text = "\n".join(c.text for c in chunks)
    for expected in ("Houston Office Under Contract", "One City Centre",
                     "Servicer commentary updated"):
        assert expected in text


def test_pptx_table_only_slide_is_indexed(tmp_path):
    """A slide made only of a table produced no chunk at all."""
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    shape = slide.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(6), Inches(2))
    t = shape.table
    t.cell(0, 0).text = "Cap rate"
    t.cell(0, 1).text = "6.5%"
    t.cell(1, 0).text = "NOI"
    t.cell(1, 1).text = "1,200,000"
    path = tmp_path / "t.pptx"
    prs.save(str(path))

    chunks = chunk_pptx(path)
    assert [c.locator for c in chunks] == ["slide 1"]
    assert "Cap rate" in chunks[0].text
    assert "1,200,000" in chunks[0].text


def test_xlsx_sheets_and_formulas(tmp_path):
    """FINC389 is a spreadsheet course: the formula is the content, so it must
    be indexed as written rather than reduced to a cached number."""
    import openpyxl

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Loan"
    ws["A1"] = "Principal"
    ws["B1"] = 250000
    ws["A2"] = "Payment"
    ws["B2"] = "=PMT(B4/12,B5*12,-B1)"
    ws2 = wb.create_sheet("Notes")
    ws2["A1"] = "Amortization schedule assumptions"
    path = tmp_path / "model.xlsx"
    wb.save(str(path))

    chunks = chunk_file(path)
    locs = [c.locator for c in chunks]
    assert locs == ["sheet Loan", "sheet Notes"]
    assert "=PMT(B4/12,B5*12,-B1)" in chunks[0].text
    assert "Principal | 250000" in chunks[0].text
    assert " |  | " not in chunks[0].text, "layout gaps must not become separators"
    assert "Amortization schedule assumptions" in chunks[1].text


def test_xlsx_empty_sheet_yields_no_chunk(tmp_path):
    import openpyxl

    wb = openpyxl.Workbook()
    wb.active.title = "Blank"
    path = tmp_path / "empty.xlsx"
    wb.save(str(path))
    assert chunk_file(path) == []


def test_html_text_only_scripts_dropped(tmp_path):
    p = tmp_path / "toc.html"
    p.write_text(
        "<html><head><title>Table of Contents</title>"
        "<style>body{color:red}</style></head><body>"
        "<script>var secret = 'do not index me';</script>"
        "<h1>Week 1</h1><p>Introduction to Banking.</p>"
        "</body></html>",
        encoding="utf-8",
    )
    chunks = chunk_file(p)
    assert chunks
    text = "\n".join(c.text for c in chunks)
    assert "Introduction to Banking." in text
    assert "Week 1" in text
    assert "do not index me" not in text
    assert "color:red" not in text
    assert chunks[0].locator.startswith("Table of Contents")


def test_utf8_bom_does_not_hide_the_first_heading(tmp_path):
    p = tmp_path / "bom.md"
    p.write_bytes(b"\xef\xbb\xbf# Alpha\nAlpha body.\n\n## Beta\nBeta body.\n")
    chunks = chunk_file(p)
    assert [c.locator for c in chunks] == ["Alpha", "Alpha > Beta"]


def test_utf16_is_decoded_not_mojibake(tmp_path):
    p = tmp_path / "u16.md"
    p.write_bytes("# Título\nEl subjuntivo es importante.\n".encode("utf-16-le"))
    chunks = chunk_file(p)
    text = "\n".join(c.text for c in chunks)
    assert "subjuntivo" in text
    assert "\x00" not in text


def test_chunks_respect_a_custom_tokenizer_and_target(tmp_path):
    """The indexer packs with the embedding model's real tokenizer; a chunk
    must not exceed the window it was told about."""
    p = tmp_path / "big.md"
    p.write_text("# H\n" + " ".join(sentence(i) for i in range(400)), encoding="utf-8")

    def count(text: str) -> int:  # 1 token per whitespace-separated word
        return len(text.split())

    chunks = chunk_file(p, count_tokens=count, target_tokens=60)
    assert len(chunks) > 1
    for c in chunks:
        assert count(c.text) <= 60, f"chunk over the window: {count(c.text)}"


def test_docx_with_oversized_attribute_falls_back(tmp_path):
    """Scanned readings embed image data as attributes larger than lxml's
    10 MB cap; those documents must still index instead of failing."""
    import zipfile

    import docx

    src = tmp_path / "src.docx"
    d = docx.Document()
    d.add_heading("Chapter One", level=1)
    d.add_paragraph("Body of chapter one.")
    d.add_heading("Section A", level=2)
    d.add_paragraph("Body of section A.")
    d.save(str(src))

    # Rewrite document.xml with an attribute value over the parser limit.
    huge = tmp_path / "huge.docx"
    with zipfile.ZipFile(src) as zin, zipfile.ZipFile(huge, "w") as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            if item.filename == "word/document.xml":
                text = data.decode("utf-8")
                blob = "A" * (11 * 1024 * 1024)
                text = text.replace("<w:body>", f'<w:body><w:p><w:pPr/><w:r><w:t data-blob="{blob}">x</w:t></w:r></w:p>', 1)
                data = text.encode("utf-8")
            zout.writestr(item, data)

    # The default python-docx path must genuinely fail on this input...
    from lxml.etree import XMLSyntaxError
    try:
        docx.Document(str(huge))
        raise AssertionError("expected lxml to reject the oversized attribute")
    except XMLSyntaxError:
        pass

    # ...and chunk_file must still produce usable chunks with heading paths.
    chunks = chunk_file(huge)
    locs = [c.locator for c in chunks]
    assert "Chapter One" in locs
    assert "Chapter One > Section A" in locs
    assert any("Body of section A." in c.text for c in chunks)


def test_txt_locators():
    small = chunk_txt("Just one small sentence.")
    assert [c.locator for c in small] == ["full text"]
    big = chunk_txt(" ".join(sentence(i) for i in range(300)))
    assert all(re.match(r"part \d+", c.locator) for c in big)


def test_unsupported_suffix_raises(tmp_path):
    p = tmp_path / "t.xyz"
    p.write_text("hello")
    try:
        chunk_file(p)
        raise AssertionError("expected ParseError")
    except ParseError as e:
        assert "Unsupported" in e.reason


def test_corrupt_pdf_raises_parse_error(tmp_path):
    p = tmp_path / "bad.pdf"
    p.write_bytes(b"this is not a pdf at all")
    try:
        chunk_file(p)
        raise AssertionError("expected ParseError")
    except ParseError as e:
        assert str(p) in str(e)
